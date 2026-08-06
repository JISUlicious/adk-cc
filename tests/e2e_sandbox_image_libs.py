"""The sandbox image's libraries, exercised inside the real container.

`import xgboost` succeeding proves almost nothing about this image: the ways
these particular packages fail are all one step past the import.

  - xgboost imports fine and then dies on `libgomp.so.1` at BOOSTER time on a
    -slim base,
  - plotly imports fine and `write_image()` fails because kaleido has no
    working binary for the arch,
  - fonts-nanum installs fine and matplotlib still renders tofu because its
    font cache never saw the new faces,
  - pyhwp is the one package here old enough to be a real casualty of the
    3.13 bump.

So every check below produces an ARTIFACT — a written file, a trained model,
rendered pixels — rather than asserting a module object exists.

Skips cleanly when Docker is unavailable, which is the normal state on a dev
box; it is the sandbox HOST that must pass this.

Run:  .venv/bin/python tests/e2e_sandbox_image_libs.py
      ADK_CC_SANDBOX_IMAGE=my-image:tag .venv/bin/python tests/...
"""
from __future__ import annotations

import os
import shutil
import subprocess
import sys

IMAGE = os.environ.get("ADK_CC_SANDBOX_IMAGE", "adk-cc-sandbox:latest")

# Runs as pid 1 inside the container. Prints one "name: OK/FAIL detail" line
# per check so a failure names itself instead of surfacing as a traceback.
PROBE = r'''
import io, os, sys, traceback, warnings
warnings.filterwarnings("ignore")
os.chdir("/tmp")
results = []

def check(name):
    def deco(fn):
        try:
            detail = fn() or ""
            results.append((name, True, str(detail)))
        except Exception as exc:
            results.append((name, False, f"{type(exc).__name__}: {exc}"))
        return fn
    return deco

@check("xlsx round-trip (openpyxl)")
def _():
    import openpyxl
    wb = openpyxl.Workbook(); wb.active["A1"] = "매출"; wb.active["B1"] = 42
    wb.save("t.xlsx")
    got = openpyxl.load_workbook("t.xlsx").active
    assert (got["A1"].value, got["B1"].value) == ("매출", 42), got["A1"].value
    return "hangul cell survived"

@check("xlsx write (xlsxwriter + chart)")
def _():
    import xlsxwriter
    wb = xlsxwriter.Workbook("t2.xlsx"); ws = wb.add_worksheet()
    ws.write_column("A1", [1, 2, 3])
    ch = wb.add_chart({"type": "line"}); ch.add_series({"values": "=Sheet1!$A$1:$A$3"})
    ws.insert_chart("C1", ch); wb.close()
    return f"{os.path.getsize('t2.xlsx')} bytes"

@check("xls read (xlrd)")
def _():
    import xlrd; return xlrd.__version__

@check("docx round-trip (python-docx)")
def _():
    import docx
    d = docx.Document(); d.add_paragraph("보고서 초안"); d.save("t.docx")
    assert docx.Document("t.docx").paragraphs[0].text == "보고서 초안"
    return "ok"

@check("pptx write (python-pptx)")
def _():
    import pptx
    p = pptx.Presentation(); s = p.slides.add_slide(p.slide_layouts[5])
    s.shapes.title.text = "분기 실적"; p.save("t.pptx")
    assert pptx.Presentation("t.pptx").slides[0].shapes.title.text == "분기 실적"
    return "ok"

@check("odf (odfpy)")
def _():
    from odf.opendocument import OpenDocumentText
    from odf.text import P
    d = OpenDocumentText(); d.text.addElement(P(text="ok")); d.save("t.odt")
    return "ok"

@check("hwpx read path (lxml + zipfile)")
def _():
    # .hwpx is a zip of XML; this is the whole reader, so prove it end to end
    # by building a minimal one and parsing its section back out.
    import zipfile
    from lxml import etree
    with zipfile.ZipFile("t.hwpx", "w") as z:
        z.writestr("Contents/section0.xml",
                   '<?xml version="1.0"?><sec><p>한글 문서</p></sec>')
    with zipfile.ZipFile("t.hwpx") as z:
        root = etree.fromstring(z.read("Contents/section0.xml"))
    assert root.findtext("p") == "한글 문서"
    return "parsed"

@check("hwp OLE container (olefile)")
def _():
    import olefile; return f"olefile {olefile.__version__}"

@check("hwp v5 reader (pyhwp) [optional]")
def _():
    import hwp5
    from hwp5.dataio import ParseError  # noqa: F401  real submodule, not a stub
    return getattr(hwp5, "__version__", "present")

@check("xgboost TRAINS (libgomp linked)")
def _():
    # The import is not the risk; building a booster is what touches OpenMP.
    import numpy as np, xgboost as xgb
    rng = np.random.default_rng(0)
    X, y = rng.normal(size=(64, 4)), rng.integers(0, 2, 64)
    m = xgb.XGBClassifier(n_estimators=8, max_depth=2, tree_method="hist")
    m.fit(X, y)
    assert m.predict(X[:4]).shape == (4,)
    return f"xgboost {xgb.__version__} trained"

@check("plotly writes a PNG (kaleido)")
def _():
    # The image pins kaleido 0.x because 1.x needs an 800MB+ browser. That pin
    # is deprecated upstream, so THIS is the check that tells us the day a
    # plotly upgrade drops v0 support — see the Dockerfile note.
    import warnings as w
    with w.catch_warnings():
        w.simplefilter("ignore", DeprecationWarning)
        import plotly.graph_objects as go
        go.Figure(go.Bar(x=["가", "나"], y=[1, 2])).write_image("t.png")
        n = os.path.getsize("t.png")
        assert n > 1000, f"suspiciously small png: {n}"
        go.Figure(go.Bar(x=[1], y=[2])).write_image("t.svg")   # vector for slides
    return f"{n} bytes png + svg"

@check("no GPU payload (xgboost-cpu, not xgboost)")
def _():
    # 244MB of CUDA libs slipped in the first time this image was built; the
    # cost is invisible from Python, so assert on the filesystem.
    import sysconfig
    nv = os.path.join(sysconfig.get_paths()["purelib"], "nvidia")
    assert not os.path.isdir(nv), f"CUDA libs present at {nv}"
    return "no nvidia/ in site-packages"

@check("matplotlib RENDERS Hangul (not tofu)")
def _():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager
    fams = sorted({f.name for f in font_manager.fontManager.ttflist
                   if f.name.startswith("Nanum")})
    assert fams, "no Nanum face visible to matplotlib"
    fig, ax = plt.subplots()
    ax.set_title("월별 매출 추이"); ax.plot([1, -2, 3])
    fig.savefig("t_ko.png", dpi=80)
    # A glyph miss is only a warning in matplotlib, never an exception, so
    # catch it the one way that is unambiguous: ask the font for the codepoint.
    from matplotlib.font_manager import findfont, FontProperties
    from matplotlib.ft2font import FT2Font
    path = findfont(FontProperties(family=["NanumGothic"]))
    assert FT2Font(path).get_char_index(ord("매")), f"no Hangul glyph in {path}"
    assert matplotlib.rcParams["axes.unicode_minus"] is False, "minus not patched"
    return f"{len(fams)} Nanum families: {', '.join(fams[:3])}"

@check("core analysis stack still intact")
def _():
    import pandas, numpy, scipy, sklearn, seaborn, polars, pyarrow, duckdb
    duckdb.sql("select 1")
    return f"pandas {pandas.__version__}, polars {polars.__version__}"

@check("runs unprivileged")
def _():
    assert os.getuid() == 1000, f"uid={os.getuid()}"
    return "uid 1000"

@check("uv is installed")
def _():
    # The gap this catches: the image shipped every analysis package while
    # missing the tool that provisions the RUNTIME interpreter, so a library
    # probe passed 15/15 and every skill script still died with "uv is not
    # available in the execution environment". Presence is checked here;
    # whether it can actually BUILD an env needs network, so that runs
    # separately (see UV_PROBE).
    import subprocess
    v = subprocess.run(["uv", "--version"], capture_output=True, text=True)
    assert v.returncode == 0, f"uv missing: {v.stderr[:200]}"
    return v.stdout.strip()

for name, ok, detail in results:
    print(f"{'OK  ' if ok else 'FAIL'}|{name}|{detail}")
print("PYTHON|" + sys.version.split()[0])
'''

# Run WITH network, unlike the probe above: this reproduces what
# `analysis_env._provision()` does inside the sandbox on the first skill run —
# `uv venv --clear --python 3.12` (which DOWNLOADS an interpreter, since the
# image ships 3.13) then `uv pip install`. Checking `uv --version` alone would
# have missed a uv that cannot actually fetch anything, which on a restricted
# network is the likely failure.
UV_PROBE = r'''
import os, subprocess, sys
os.chdir("/tmp")
steps = [
    (["uv", "venv", "--clear", "--python", "3.12", "/tmp/probe-env"],
     "uv venv --python 3.12 (downloads an interpreter)"),
    (["uv", "pip", "install", "--quiet", "--python",
      "/tmp/probe-env/bin/python", "packaging"], "uv pip install"),
    (["/tmp/probe-env/bin/python", "-c",
      "import sys, packaging; print(sys.version.split()[0], packaging.__version__)"],
     "the provisioned interpreter runs with the installed package"),
]
for cmd, what in steps:
    p = subprocess.run(cmd, capture_output=True, text=True, timeout=900)
    if p.returncode != 0:
        print(f"FAIL|{what}|{(p.stderr or p.stdout).strip()[-300:]}")
        sys.exit(0)
    last = (p.stdout or "").strip()
print(f"OK  |uv provisions an env exactly as analysis_env does|{last}")
'''


def main() -> int:
    if not shutil.which("docker"):
        print("SKIP: docker not installed."); return 0
    if subprocess.run(["docker", "info"], capture_output=True).returncode != 0:
        print("SKIP: docker daemon not running."); return 0
    has = subprocess.run(["docker", "image", "inspect", IMAGE],
                         capture_output=True)
    if has.returncode != 0:
        print(f"SKIP: {IMAGE} not built "
              f"(docker build -t {IMAGE} -f Dockerfile.sandbox .)"); return 0

    print(f"probing {IMAGE}")
    # --network none on purpose: every library check must work offline, which
    # is a real property (kaleido 0.2.1 renders without fetching anything).
    proc = subprocess.run(
        ["docker", "run", "--rm", "-i", "--network", "none", IMAGE,
         "python", "-c", PROBE],
        capture_output=True, text=True, timeout=600)
    if proc.returncode != 0 and not proc.stdout.strip():
        print("FAIL: probe did not run\n" + (proc.stderr or "")[-2000:])
        return 1

    # Provisioning genuinely needs egress, so it gets its own networked run
    # rather than weakening the isolation above.
    uv = subprocess.run(
        ["docker", "run", "--rm", "-i", IMAGE, "python", "-c", UV_PROBE],
        capture_output=True, text=True, timeout=1200)
    uv_lines = [ln for ln in (uv.stdout or "").splitlines() if "|" in ln]
    if not uv_lines:
        uv_lines = [f"FAIL|uv provisioning probe did not run|"
                    f"{(uv.stderr or '').strip()[-300:]}"]

    failed = optional_failed = 0
    for line in proc.stdout.splitlines() + uv_lines:
        if line.startswith("PYTHON|"):
            print(f"  python: {line.split('|', 1)[1]}")
            continue
        if "|" not in line:
            continue
        status, name, detail = line.split("|", 2)
        ok = status.strip() == "OK"
        # pyhwp is allowed to be absent — the Dockerfile installs it in a
        # non-fatal layer for exactly this reason. Report it, don't fail on it.
        optional = "[optional]" in name
        if not ok and optional:
            optional_failed += 1
        elif not ok:
            failed += 1
        mark = "PASS" if ok else ("WARN" if optional else "FAIL")
        print(f"  [{mark}] {name}" + (f" — {detail}" if detail else ""))

    print(f"\n{failed} required checks failed"
          + (f", {optional_failed} optional missing" if optional_failed else ""))
    if optional_failed:
        print("  NOTE: binary .hwp reads are unavailable in this image; "
              ".hwpx still works.")
    return 1 if failed else 0


if __name__ == "__main__":
    sys.exit(main())
